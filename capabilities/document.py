import os
import re
from pathlib import Path
from typing import Any, Dict
from dataclasses import dataclass
from loguru import logger

# Document libraries
try:
    from docx import Document
    from docx.shared import Pt
    from docx.enum.text import WD_PARAGRAPH_ALIGNMENT
    from docx.shared import RGBColor as DocxRGBColor
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor as PptxRGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from pypdf import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.units import inch
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False


@dataclass
class CapabilityResult:
    success: bool
    output: str
    error: str | None = None


class DocumentCapability:
    """Handles Word docs, PowerPoint, PDF, Markdown, and text files."""

    def __init__(self):
        self._check_dependencies()

    def _check_dependencies(self):
        """Check if required libraries are available."""
        missing = []
        if not DOCX_AVAILABLE:
            missing.append("python-docx")
        if not PPTX_AVAILABLE:
            missing.append("python-pptx")
        if not PDF_AVAILABLE:
            missing.append("pypdf")
        if not REPORTLAB_AVAILABLE:
            missing.append("reportlab")
        
        if missing:
            logger.warning(
                f"Document capability limited - missing: {', '.join(missing)}"
            )

    def _resolve(self, path: str) -> Path:
        """Resolve path handling ~, %USERPROFILE%, absolute and relative."""
        path = path.strip()
        
        # Handle shell expansion
        if path.startswith("~"):
            path = os.path.expanduser(path)
        elif "%USERPROFILE%" in path.upper():
            path = path.replace("%USERPROFILE%", os.environ.get("USERPROFILE", ""))
            path = path.replace("%userprofile%", os.environ.get("USERPROFILE", ""))
        
        # Convert to absolute
        p = Path(path)
        if not p.is_absolute():
            p = Path.cwd() / p
        
        return p.resolve()

    def execute(self, operation: str, params: Dict[str, Any]) -> CapabilityResult:
        """Execute a document operation."""
        try:
            if operation == "create_word_doc":
                return self._create_word_doc(params)
            elif operation == "edit_word_doc":
                return self._edit_word_doc(params)
            elif operation == "read_word_doc":
                return self._read_word_doc(params)
            elif operation == "create_powerpoint":
                return self._create_powerpoint(params)
            elif operation == "edit_powerpoint":
                return self._edit_powerpoint(params)
            elif operation == "read_powerpoint":
                return self._read_powerpoint(params)
            elif operation == "create_pdf":
                return self._create_pdf(params)
            elif operation == "read_pdf":
                return self._read_pdf(params)
            elif operation == "create_markdown":
                return self._create_markdown(params)
            elif operation == "edit_markdown":
                return self._edit_markdown(params)
            elif operation == "read_text":
                return self._read_text(params)
            elif operation == "edit_text":
                return self._edit_text(params)
            else:
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"Unknown document operation: {operation}"
                )
        except Exception as e:
            logger.error(f"Document operation failed: {e}")
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _create_word_doc(self, params: Dict[str, Any]) -> CapabilityResult:
        """Create a Word document."""
        if not DOCX_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="python-docx not available"
            )

        path = self._resolve(params["path"])
        title = params.get("title", "")
        content = params.get("content", "")

        try:
            doc = Document()
            
            # Add title if provided
            if title:
                title_para = doc.add_heading(title, level=1)
                title_run = title_para.runs[0]
                title_run.font.color.rgb = DocxRGBColor(0, 0, 0)
                title_run.font.size = Pt(16)

            # Parse content for headings and paragraphs
            lines = content.split('\n')
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # Check for headings (#, ##, ###)
                if line.startswith('#'):
                    level = len(line) - len(line.lstrip('#'))
                    level = min(level, 3)  # Max level 3
                    heading_text = line.lstrip('#').strip()
                    doc.add_heading(heading_text, level=level)
                else:
                    doc.add_paragraph(line)

            # Create parent directory if needed
            path.parent.mkdir(parents=True, exist_ok=True)
            doc.save(path)
            
            return CapabilityResult(
                success=True,
                output=f"Word document created: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _edit_word_doc(self, params: Dict[str, Any]) -> CapabilityResult:
        """Edit an existing Word document."""
        if not DOCX_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="python-docx not available"
            )

        path = self._resolve(params["path"])
        instructions = params.get("instructions", "")

        try:
            if not path.exists():
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"File not found: {path}"
                )

            doc = Document(path)
            
            # Simple implementation: append new content
            # In a more sophisticated version, we'd parse instructions
            # to find specific sections to replace
            doc.add_paragraph()  # Add space
            doc.add_paragraph(instructions)
            
            doc.save(path)
            
            return CapabilityResult(
                success=True,
                output=f"Word document edited: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _read_word_doc(self, params: Dict[str, Any]) -> CapabilityResult:
        """Read a Word document."""
        if not DOCX_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="python-docx not available"
            )

        path = self._resolve(params["path"])

        try:
            if not path.exists():
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"File not found: {path}"
                )

            doc = Document(path)
            content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            text_content = '\n'.join(content)
            
            return CapabilityResult(
                success=True,
                output=text_content
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _create_powerpoint(self, params: Dict[str, Any]) -> CapabilityResult:
        """Create a PowerPoint presentation."""
        if not PPTX_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="python-pptx not available"
            )

        path = self._resolve(params["path"])
        title = params.get("title", "Presentation")
        slides = params.get("slides", [])

        try:
            prs = Presentation()
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]  # Title slide layout
            slide = prs.slides.add_slide(title_slide_layout)
            title_shape = slide.shapes.title
            title_shape.text = title
            
            # Set dark theme colors
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = PptxRGBColor(26, 26, 46)  # #1a1a2e
            
            # Content slides
            for slide_data in slides:
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)
                
                slide_title = slide_data.get("title", "")
                slide_content = slide_data.get("content", "")
                
                if slide.shapes.title:
                    slide.shapes.title.text = slide_title
                    # Set title color
                    title_shape = slide.shapes.title
                    for paragraph in title_shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = PptxRGBColor(155, 89, 255)  # #9B59FF
                
                # Set content
                if slide.placeholders:
                    content_placeholder = None
                    for placeholder in slide.placeholders:
                        if placeholder.has_text_frame:
                            content_placeholder = placeholder
                            break
                    
                    if content_placeholder:
                        content_placeholder.text = slide_content
                        # Set content color to white
                        for paragraph in content_placeholder.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = PptxRGBColor(255, 255, 255)
                
                # Set background
                background = slide.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = PptxRGBColor(26, 26, 46)  # #1a1a2e

            # Create parent directory if needed
            path.parent.mkdir(parents=True, exist_ok=True)
            prs.save(path)
            
            return CapabilityResult(
                success=True,
                output=f"PowerPoint created: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _edit_powerpoint(self, params: Dict[str, Any]) -> CapabilityResult:
        """Edit a specific PowerPoint slide."""
        if not PPTX_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="python-pptx not available"
            )

        path = self._resolve(params["path"])
        slide_index = params.get("slide_index", 0)
        title = params.get("title", "")
        content = params.get("content", "")

        try:
            if not path.exists():
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"File not found: {path}"
                )

            prs = Presentation(path)
            
            if slide_index >= len(prs.slides):
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"Slide index {slide_index} out of range (0-{len(prs.slides)-1})"
                )

            slide = prs.slides[slide_index]
            
            if slide.shapes.title and title:
                slide.shapes.title.text = title
            
            if slide.placeholders and content:
                for placeholder in slide.placeholders:
                    if placeholder.has_text_frame:
                        placeholder.text = content
                        break

            prs.save(path)
            
            return CapabilityResult(
                success=True,
                output=f"PowerPoint slide {slide_index} edited: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _read_powerpoint(self, params: Dict[str, Any]) -> CapabilityResult:
        """Read PowerPoint content."""
        if not PPTX_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="python-pptx not available"
            )

        path = self._resolve(params["path"])

        try:
            if not path.exists():
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"File not found: {path}"
                )

            prs = Presentation(path)
            content = []
            
            for i, slide in enumerate(prs.slides):
                slide_content = f"Slide {i+1}:\n"
                
                if slide.shapes.title:
                    slide_content += f"Title: {slide.shapes.title.text}\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += f"Content: {shape.text}\n"
                
                content.append(slide_content)
            
            return CapabilityResult(
                success=True,
                output='\n'.join(content)
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _create_pdf(self, params: Dict[str, Any]) -> CapabilityResult:
        """Create a PDF using reportlab."""
        if not REPORTLAB_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="reportlab not available"
            )

        path = self._resolve(params["path"])
        title = params.get("title", "")
        content = params.get("content", "")

        try:
            # Create parent directory if needed
            path.parent.mkdir(parents=True, exist_ok=True)
            
            c = canvas.Canvas(str(path), pagesize=letter)
            width, height = letter
            
            # Title
            if title:
                c.setFont("Helvetica-Bold", 16)
                c.drawString(inch, height - inch, title)
                y_position = height - 1.5 * inch
            else:
                y_position = height - inch
            
            # Content
            c.setFont("Helvetica", 12)
            lines = content.split('\n')
            
            for line in lines:
                if y_position < inch:  # Start new page if needed
                    c.showPage()
                    y_position = height - inch
                
                c.drawString(inch, y_position, line)
                y_position -= 0.2 * inch
            
            c.save()
            
            return CapabilityResult(
                success=True,
                output=f"PDF created: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _read_pdf(self, params: Dict[str, Any]) -> CapabilityResult:
        """Read PDF content using pypdf."""
        if not PDF_AVAILABLE:
            return CapabilityResult(
                success=False,
                output="",
                error="pypdf not available"
            )

        path = self._resolve(params["path"])

        try:
            if not path.exists():
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"File not found: {path}"
                )

            reader = PdfReader(path)
            content = []
            
            for page in reader.pages:
                page_text = page.extract_text()
                if page_text.strip():
                    content.append(page_text)
            
            return CapabilityResult(
                success=True,
                output='\n'.join(content)
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _create_markdown(self, params: Dict[str, Any]) -> CapabilityResult:
        """Create a Markdown file."""
        path = self._resolve(params["path"])
        content = params.get("content", "")

        try:
            # Create parent directory if needed
            path.parent.mkdir(parents=True, exist_ok=True)
            
            path.write_text(content, encoding='utf-8')
            
            return CapabilityResult(
                success=True,
                output=f"Markdown file created: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _edit_markdown(self, params: Dict[str, Any]) -> CapabilityResult:
        """Edit a Markdown file."""
        path = self._resolve(params["path"])
        content = params.get("content", "")
        mode = params.get("mode", "append")

        try:
            if mode == "overwrite":
                path.write_text(content, encoding='utf-8')
            else:  # append
                if path.exists():
                    existing = path.read_text(encoding='utf-8')
                    path.write_text(existing + '\n' + content, encoding='utf-8')
                else:
                    path.write_text(content, encoding='utf-8')
            
            return CapabilityResult(
                success=True,
                output=f"Markdown file edited: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _read_text(self, params: Dict[str, Any]) -> CapabilityResult:
        """Read any text-based file."""
        path = self._resolve(params["path"])

        try:
            if not path.exists():
                return CapabilityResult(
                    success=False,
                    output="",
                    error=f"File not found: {path}"
                )

            content = path.read_text(encoding='utf-8')
            
            return CapabilityResult(
                success=True,
                output=content
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )

    def _edit_text(self, params: Dict[str, Any]) -> CapabilityResult:
        """Edit any text-based file."""
        path = self._resolve(params["path"])
        content = params.get("content", "")
        mode = params.get("mode", "append")

        try:
            # Create parent directory if needed
            path.parent.mkdir(parents=True, exist_ok=True)
            
            if mode == "overwrite":
                path.write_text(content, encoding='utf-8')
            else:  # append
                if path.exists():
                    existing = path.read_text(encoding='utf-8')
                    path.write_text(existing + '\n' + content, encoding='utf-8')
                else:
                    path.write_text(content, encoding='utf-8')
            
            return CapabilityResult(
                success=True,
                output=f"Text file edited: {path}"
            )
        except Exception as e:
            return CapabilityResult(
                success=False,
                output="",
                error=str(e)
            )
